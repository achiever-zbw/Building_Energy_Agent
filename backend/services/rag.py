""" LightRAG 服务（HTTP 与 MCP 子进程共用配置；数据目录由 WORKING_DIR 对齐）。"""
from __future__ import annotations

import asyncio
import io
import os
from pathlib import Path
from typing import Any

import numpy as np
from docx import Document
from fastapi import HTTPException, UploadFile
from lightrag import LightRAG
from lightrag.kg.shared_storage import initialize_pipeline_status
from lightrag.llm.openai import openai_complete_if_cache, openai_embed
from lightrag.utils import EmbeddingFunc
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

_mcp_service_lock = asyncio.Lock()
_mcp_lightrag_service: LightRAGService | None = None

# 加载环境变量
def _load_env() -> None:
    try:
        from dotenv import load_dotenv
    except ImportError:
        return
    backend_root = Path(__file__).resolve().parent.parent
    project_root = backend_root.parent
    load_dotenv(project_root / ".env", override=False)
    load_dotenv(backend_root / ".env", override=False)


def _working_dir() -> str:
    """
    LightRAG 数据目录。相对路径一律相对「项目根」（backend 的上一级），
    避免主进程 cwd 与 MCP 子进程 cwd=backend 时解析成两套目录。
    """
    _load_env()
    backend_root = Path(__file__).resolve().parent.parent
    project_root = backend_root.parent
    raw = (os.getenv("WORKING_DIR") or "").strip()
    if raw:
        p = Path(raw)
        if not p.is_absolute():
            p = (project_root / p).resolve()
        else:
            p = p.resolve()
        return str(p)
    return str((backend_root / "data" / "lightrag").resolve())


async def llm_model_func(
    prompt,
    system_prompt=None,
    history_messages=None,
    keyword_extraction=False,
    **kwargs,
) -> str:
    if history_messages is None:
        history_messages = []
    return await openai_complete_if_cache(
        model=os.getenv("MODEL"),
        prompt=prompt,
        system_prompt=system_prompt,
        history_messages=history_messages,
        api_key=os.getenv("OPENAI_API_KEY"),
        base_url=os.getenv("BASE_URL"),
        **kwargs,
    )


async def embedding_func(texts: list[str]) -> np.ndarray:
    embed_fn = getattr(openai_embed, "func", openai_embed)
    return await embed_fn(
        texts,
        model=os.getenv("EMBEDDING_MODEL"),
        api_key=os.getenv("EMBEDDING_API_KEY"),
        base_url=os.getenv("EMBEDDING_MODEL_URL"),
    )


class LightRAGService:
    """ LightRAG 服务 """

    def __init__(self) -> None:
        self.rag: LightRAG | None = None

    async def initialize_lightrag(self) -> LightRAG:
        _load_env()
        wd = _working_dir()
        Path(wd).mkdir(parents=True, exist_ok=True)
        dim = int(os.getenv("EMBEDDING_DIM", "1536"))

        self.rag = LightRAG(
            working_dir=wd,
            embedding_func=EmbeddingFunc(
                embedding_dim=dim,
                func=embedding_func,
            ),
            llm_model_func=llm_model_func,
        )

        await self.rag.initialize_storages()
        await initialize_pipeline_status()
        return self.rag

    def _require_rag(self) -> LightRAG:
        if self.rag is None:
            raise RuntimeError("LightRAG 未初始化")
        return self.rag

    async def rag_query(self, query: str) -> Any:
        """查询；供 HTTP/MCP 使用。同步 query 放到线程避免阻塞事件循环。"""
        if not (query or "").strip():
            return {"error": "问题为空"}
        try:
            rag = self._require_rag()
        except RuntimeError as e:
            return {"error": str(e)}

        def _run() -> Any:
            try:
                return rag.query(query=query)
            except TypeError:
                return rag.query(query)

        return await asyncio.to_thread(_run)

    async def _extract_text_by_extension(self, raw: bytes, ext: str) -> str:
        text_exts = {
            ".txt", ".md", ".html", ".htm", ".tex", ".json", ".xml",
            ".yaml", ".yml", ".rtf", ".odt", ".epub", ".csv",
            ".log", ".conf", ".ini", ".properties", ".sql",
            ".bat", ".sh", ".c", ".cpp", ".py", ".java", ".js",
            ".ts", ".swift", ".go", ".rb", ".php", ".css", ".scss", ".less",
        }
        if ext in text_exts:
            return raw.decode("utf-8")
        if ext == ".pdf":
            def _pdf_extract() -> str:
                reader = PdfReader(io.BytesIO(raw))
                return "\n".join((p.extract_text() or "") for p in reader.pages)
            return await asyncio.to_thread(_pdf_extract)
        if ext == ".docx":
            def _docx_extract() -> str:
                doc = Document(io.BytesIO(raw))
                return "\n".join(p.text for p in doc.paragraphs)
            return await asyncio.to_thread(_docx_extract)
        if ext == ".pptx":
            def _pptx_extract() -> str:
                prs = Presentation(io.BytesIO(raw))
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            lines.append(shape.text)
                return "\n".join(lines)
            return await asyncio.to_thread(_pptx_extract)
        if ext == ".xlsx":
            def _xlsx_extract() -> str:
                wb = load_workbook(io.BytesIO(raw), data_only=True)
                lines = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(v) for v in row if v is not None]
                        if vals:
                            lines.append(" | ".join(vals))
                return "\n".join(lines)
            return await asyncio.to_thread(_xlsx_extract)
        raise ValueError(f"Unsupported file type: {ext}")

    async def upload_document(self, document_files: list[UploadFile]) -> dict[str, Any]:
        if self.rag is None:
            raise HTTPException(status_code=503, detail="LightRAG 未初始化")
        rag = self.rag

        if not document_files:
            raise HTTPException(status_code=400, detail="No files provided")

        supported_extensions = {
            ".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx",
            ".html", ".htm", ".json", ".xml", ".yaml", ".yml", ".csv",
        }

        success_files: list[str] = []
        failed_files: list[dict[str, Any]] = []

        for upload_file in document_files:
            filename = (upload_file.filename or "").strip()
            ext = Path(filename).suffix.lower()

            if not filename:
                failed_files.append({"file": filename, "error": "empty filename"})
                continue

            if ext not in supported_extensions:
                failed_files.append({"file": filename, "error": f"unsupported type: {ext}"})
                continue

            try:
                raw = await upload_file.read()
                if not raw:
                    failed_files.append({"file": filename, "error": "empty file"})
                    continue

                content = await self._extract_text_by_extension(raw=raw, ext=ext)

                if not content or not content.strip():
                    failed_files.append({"file": filename, "error": "no extractable text"})
                    continue

                await rag.apipeline_enqueue_documents(
                    input=content,
                    file_paths=filename,
                    track_id=None,
                )
                success_files.append(filename)

            except Exception as e:
                failed_files.append({"file": filename, "error": str(e)})
            finally:
                await upload_file.close()

        if success_files:
            await rag.apipeline_process_enqueue_documents()

        return {
            "success_count": len(success_files),
            "failed_count": len(failed_files),
            "success_files": success_files,
            "failed_files": failed_files,
        }


async def get_lightrag_service_for_mcp() -> LightRAGService:
    """
    MCP 独立子进程内使用：与 FastAPI 进程分离，不能读 app.state。
    同一 WORKING_DIR 下与主进程共享落盘数据。
    """
    global _mcp_lightrag_service
    async with _mcp_service_lock:
        if _mcp_lightrag_service is None:
            svc = LightRAGService()
            await svc.initialize_lightrag()
            _mcp_lightrag_service = svc
        return _mcp_lightrag_service
